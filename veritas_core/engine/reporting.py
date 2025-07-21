# ==============================================================================
# Core Engine: Compliant Report Generation
#
# Author: Principal Engineer SME
# Last Updated: 2025-07-20
#
# Description:
# This module is the centralized engine for generating GxP-compliant,
# submission-ready reports in PDF and PowerPoint formats for the VERITAS
# application. It supports dynamic content assembly, draft watermarking, and
# electronic signatures, ensuring robustness, flexibility, and adherence to data
# integrity principles.
#
# Key Features:
# - Dynamic Content Assembly: Builds reports from a configuration dictionary.
# - GxP Watermarking: Adds "DRAFT" watermark for review copies in PDFs.
# - E-Signature Placeholder: Includes a formatted section for electronic signatures.
# ==============================================================================

import io
import pandas as pd
from datetime import datetime
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import plotly.graph_objects as go
from typing import Dict, Any

# --- PDF Generation Engine ---

class VeritasPDF(FPDF):
    """
    Custom FPDF class with VERITAS branding, headers, footers, and watermarking.
    """
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.watermark_text = ""

    def set_watermark(self, text: str) -> None:
        """
        Set the watermark text for the PDF.

        Args:
            text (str): Watermark text to display.

        Raises:
            ValueError: If text is not a non-empty string.
        """
        if not isinstance(text, str) or not text.strip():
            raise ValueError("Watermark text must be a non-empty string")
        self.watermark_text = text

    def header(self) -> None:
        """
        Add a header with VERITAS branding and generation timestamp.

        Raises:
            RuntimeError: If font setting or cell rendering fails.
        """
        try:
            self.set_font('Helvetica', 'B', 12)
            self.cell(0, 10, 'VERITAS - Automated Data Summary Report', 0, 1, 'C')
            self.set_font('Helvetica', '', 8)
            self.cell(0, 5, f"Generated: {datetime.now(tz=timezone.utc).strftime('%Y-%m-%d %H:%M:%S UTC')}", 0, 1, 'C')
            self.ln(10)
        except Exception as e:
            raise RuntimeError(f"Failed to render PDF header: {str(e)}")

    def footer(self) -> None:
        """
        Add a footer with page number and confidentiality notice.

        Raises:
            RuntimeError: If font setting or cell rendering fails.
        """
        try:
            self.set_y(-15)
            self.set_font('Helvetica', 'I', 8)
            self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
            self.set_x(-50)
            self.cell(0, 10, 'VERTEX CONFIDENTIAL', 0, 0, 'R')
        except Exception as e:
            raise RuntimeError(f"Failed to render PDF footer: {str(e)}")

    def _add_watermark(self) -> None:
        """
        Add a watermark to the current page if set.

        Raises:
            RuntimeError: If watermark rendering fails.
        """
        if self.watermark_text:
            try:
                self.set_font('Helvetica', 'B', 50)
                self.set_text_color(220, 220, 220)
                self.rotate(45)
                self.text(35, 190, self.watermark_text)
                self.rotate(0)
                self.set_text_color(0, 0, 0)
            except Exception as e:
                raise RuntimeError(f"Failed to render watermark: {str(e)}")

    def add_page(self, orientation: str = '') -> None:
        """
        Add a new page and apply watermark.

        Args:
            orientation (str): Page orientation ('P' or 'L'). Defaults to ''.

        Raises:
            RuntimeError: If page addition fails.
        """
        try:
            super().add_page(orientation)
            self._add_watermark()
        except Exception as e:
            raise RuntimeError(f"Failed to add PDF page: {str(e)}")

    def chapter_title(self, title: str) -> None:
        """
        Add a chapter title with formatting.

        Args:
            title (str): Title text.

        Raises:
            ValueError: If title is not a non-empty string.
            RuntimeError: If title rendering fails.
        """
        if not isinstance(title, str) or not title.strip():
            raise ValueError("Title must be a non-empty string")
        try:
            self.set_font('Helvetica', 'B', 12)
            self.set_fill_color(230, 230, 230)
            self.cell(0, 8, title, 0, 1, 'L', fill=True)
            self.ln(4)
        except Exception as e:
            raise RuntimeError(f"Failed to render chapter title: {str(e)}")

    def chapter_body(self, body: str) -> None:
        """
        Add a chapter body with text.

        Args:
            body (str): Body text.

        Raises:
            ValueError: If body is not a non-empty string.
            RuntimeError: If body rendering fails.
        """
        if not isinstance(body, str) or not body.strip():
            raise ValueError("Body text must be a non-empty string")
        try:
            self.set_font('Helvetica', '', 10)
            self.multi_cell(0, 5, body)
            self.ln()
        except Exception as e:
            raise RuntimeError(f"Failed to render chapter body: {str(e)}")

    def add_dataframe(self, df: pd.DataFrame, title: str) -> None:
        """
        Add a DataFrame as a table to the PDF.

        Args:
            df (pd.DataFrame): DataFrame to render.
            title (str): Table title.

        Raises:
            ValueError: If df is not a DataFrame, is empty, or title is invalid.
            RuntimeError: If table rendering fails.
        """
        if not isinstance(df, pd.DataFrame):
            raise ValueError("df must be a pandas DataFrame")
        if not isinstance(title, str) or not title.strip():
            raise ValueError("Title must be a non-empty string")
        if df.empty:
            return  # Skip rendering for empty DataFrame

        try:
            self.chapter_title(title)
            self.set_font('Helvetica', 'B', 9)
            effective_width = self.w - 2 * self.l_margin
            num_cols = len(df.columns)
            
            # Calculate column widths, capping at max width to avoid overflow
            col_widths = {col: min(df[col].astype(str).str.len().max(), 50) for col in df.columns}
            total_chars = sum(col_widths.values()) or 1  # Avoid division by zero
            col_width_props = {col: (width / total_chars) * effective_width for col, width in col_widths.items()}

            for col in df.columns:
                self.cell(col_width_props[col], 7, col, 1, 0, 'C')
            self.ln()
            
            self.set_font('Helvetica', '', 8)
            for _, row in df.iterrows():
                for col in df.columns:
                    value = str(row[col]) if pd.notna(row[col]) else ""
                    self.cell(col_width_props[col], 6, value, 1, 0, 'L')
                self.ln()
            self.ln(5)
        except Exception as e:
            raise RuntimeError(f"Failed to render DataFrame table: {str(e)}")

    def add_signature_section(self, signature_details: Dict) -> None:
        """
        Add an electronic signature section to the PDF.

        Args:
        signature_details (Dict): Dictionary with 'user', 'timestamp', 'reason' keys.

        Raises:
            ValueError: If signature_details is missing required keys or contains invalid values.
            RuntimeError: If signature section rendering fails.
        """
        if not isinstance(signature_details, dict):
            raise ValueError("signature_details must be a dictionary")
        required_keys = ['user', 'timestamp', 'reason']
        if not all(key in signature_details for key in required_keys):
            raise ValueError(f"signature_details must contain keys: {required_keys}")
        if not all(isinstance(signature_details[key], str) and signature_details[key].strip() for key in required_keys):
            raise ValueError("All signature_details values must be non-empty strings")

        try:
            self.chapter_title("4.0 Electronic Signature")
            sig_body = (
                f"This document was electronically signed and locked in the VERITAS system, "
                f"in accordance with 21 CFR Part 11 requirements.\n\n"
                f"**Signed By:** {signature_details['user']}\n"
                f"**Signature Timestamp:** {signature_details['timestamp']}\n"
                f"**Meaning of Signature:** {signature_details['reason']}"
            )
            self.chapter_body(sig_body)
        except Exception as e:
            raise RuntimeError(f"Failed to render signature section: {str(e)}")

def generate_pdf_report(report_data: Dict[str, Any], watermark: str = "") -> bytes:
    """
    Create a formatted PDF report from the provided data.

    Args:
        report_data (Dict[str, Any]): Dictionary with 'study_id', 'commentary',
            'sections_config', 'data', 'cqa', and optional 'signature_details' and 'plot_fig'.
        watermark (str): Optional watermark text (e.g., "DRAFT"). Defaults to "".

    Returns:
        bytes: Encoded PDF content.

    Raises:
        ValueError: If report_data is missing required keys or contains invalid values.
        TypeError: If report_data is not a dictionary or inputs are invalid.
        RuntimeError: If PDF generation fails.
    """
    if not isinstance(report_data, dict):
        raise TypeError("report_data must be a dictionary")
    required_keys = ['study_id', 'commentary', 'sections_config', 'data', 'cqa']
    if not all(key in report_data for key in required_keys):
        raise ValueError(f"report_data must contain keys: {required_keys}")
    if not all(isinstance(report_data[key], str) and report_data[key].strip() for key in ['study_id', 'commentary', 'cqa']):
        raise ValueError("study_id, commentary, and cqa must be non-empty strings")
    if not isinstance(report_data['data'], pd.DataFrame):
        raise TypeError("report_data['data'] must be a pandas DataFrame")
    if not isinstance(report_data['sections_config'], dict):
        raise TypeError("report_data['sections_config'] must be a dictionary")
    if not all(key in report_data['sections_config'] for key in ['include_summary_stats', 'include_full_dataset']):
        raise ValueError("sections_config must contain 'include_summary_stats' and 'include_full_dataset' keys")
    if not isinstance(watermark, str):
        raise TypeError("watermark must be a string")
    if report_data['cqa'] not in report_data['data'].columns:
        raise ValueError(f"CQA {report_data['cqa']} not found in DataFrame columns")

    try:
        pdf = VeritasPDF()
        if watermark:
            pdf.set_watermark(watermark)
        pdf.add_page()

        # Section 1: Metadata
        pdf.chapter_title(f"1.0 Summary for Study: {report_data['study_id']}")
        pdf.chapter_body(f"**Analyst Commentary:** {report_data['commentary']}")

        # Section 2: Data & Analysis
        pdf.chapter_title("2.0 Data & Analysis")
        if report_data['sections_config']['include_summary_stats']:
            summary_stats = report_data['data'][report_data['cqa']].describe().round(3).reset_index()
            summary_stats.columns = ['Statistic', 'Value']
            pdf.add_dataframe(summary_stats, "Summary Statistics")

        # Section 3: Full Dataset
        if report_data['sections_config']['include_full_dataset']:
            pdf.add_page()
            pdf.chapter_title("3.0 Appendix: Full Dataset")
            pdf.add_dataframe(report_data['data'], "Full Dataset")

        # Section 4: Signature
        if report_data.get('signature_details'):
            pdf.add_signature_section(report_data['signature_details'])

        return pdf.output(dest='S').encode('latin-1')
    except Exception as e:
        raise RuntimeError(f"Failed to generate PDF report: {str(e)}")

# --- PowerPoint Generation Engine ---

def _add_table_to_slide(slide: Any, df: pd.DataFrame, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
    """
    Add a pandas DataFrame as a table to a PowerPoint slide.

    Args:
        slide (Any): PowerPoint slide object.
        df (pd.DataFrame): DataFrame to render as a table.
        left (Inches): Left position of the table.
        top (Inches): Top position of the table.
        width (Inches): Table width.
        height (Inches): Table height.

    Raises:
        ValueError: If df is not a DataFrame, is empty, or dimensions are invalid.
        TypeError: If slide is not a valid slide object or dimensions are not Inches.
        RuntimeError: If table creation fails.
    """
    if not isinstance(df, pd.DataFrame):
        raise TypeError("df must be a pandas DataFrame")
    if df.empty:
        return  # Skip rendering for empty DataFrame
    if not all(isinstance(x, Inches) for x in [left, top, width, height]):
        raise TypeError("left, top, width, and height must be Inches")
    if any(x.inches <= 0 for x in [width, height]):
        raise ValueError("width and height must be positive")

    try:
        rows, cols = df.shape
        rows += 1  # Include header row
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        for i in range(cols):
            table.columns[i].width = Inches(width.inches / cols)
        for c, col in enumerate(df.columns):
            table.cell(0, c).text = str(col)
        for r, row in enumerate(df.itertuples(index=False)):
            for c, val in enumerate(row):
                table.cell(r + 1, c).text = str(val) if pd.notna(val) else ""
    except Exception as e:
        raise RuntimeError(f"Failed to add table to slide: {str(e)}")

def generate_ppt_report(report_data: Dict[str, Any]) -> bytes:
    """
    Generate a PowerPoint report with data and plots.

    Args:
        report_data (Dict[str, Any]): Dictionary with 'study_id', 'sections_config',
            'data', 'cqa', and optional 'plot_fig'.

    Returns:
        bytes: PowerPoint file content.

    Raises:
        ValueError: If report_data is missing required keys or contains invalid values.
        TypeError: If report_data is not a dictionary or plot_fig is not a Plotly Figure.
        RuntimeError: If PowerPoint generation fails.
    """
    if not isinstance(report_data, dict):
        raise TypeError("report_data must be a dictionary")
    required_keys = ['study_id', 'sections_config', 'data', 'cqa']
    if not all(key in report_data for key in required_keys):
        raise ValueError(f"report_data must contain keys: {required_keys}")
    if not all(isinstance(report_data[key], str) and report_data[key].strip() for key in ['study_id', 'cqa']):
        raise ValueError("study_id and cqa must be non-empty strings")
    if not isinstance(report_data['data'], pd.DataFrame):
        raise TypeError("report_data['data'] must be a pandas DataFrame")
    if not isinstance(report_data['sections_config'], dict):
        raise TypeError("report_data['sections_config'] must be a dictionary")
    if 'include_summary_stats' not in report_data['sections_config']:
        raise ValueError("sections_config must contain 'include_summary_stats' key")
    if report_data['cqa'] not in report_data['data'].columns:
        raise ValueError(f"CQA {report_data['cqa']} not found in DataFrame columns")
    if 'plot_fig' in report_data and not isinstance(report_data['plot_fig'], go.Figure):
        raise TypeError("report_data['plot_fig'] must be a Plotly go.Figure")

    try:
        prs = Presentation()

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "VERITAS Automated Study Report"
        slide.placeholders[1].text = (
            f"Study: {report_data['study_id']}\n"
            f"Generated: {datetime.now(tz=timezone.utc).strftime('%Y-%m-%d UTC')}"
        )

        # Slide 2: Data Summary
        if report_data['sections_config']['include_summary_stats']:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "Summary Statistics"
            summary_stats = report_data['data'][report_data['cqa']].describe().round(3).reset_index()
            summary_stats.columns = ['Statistic', 'Value']
            _add_table_to_slide(slide, summary_stats, Inches(1), Inches(1.5), Inches(8), Inches(3))

        # Slide 3: Plot
        if report_data.get('plot_fig'):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = report_data['plot_fig'].layout.title.text or "Analysis Chart"
            img_stream = io.BytesIO()
            report_data['plot_fig'].write_image(img_stream, format="png", width=800, height=450, scale=2)
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))

        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io.getvalue()
    except Exception as e:
        raise RuntimeError(f"Failed to generate PowerPoint report: {str(e)}")
